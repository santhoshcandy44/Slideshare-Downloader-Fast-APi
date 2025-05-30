import json
import os
import tempfile
from datetime import datetime
from bs4 import BeautifulSoup
from urllib.parse import urlparse
from PIL import Image
import requests
from io import BytesIO

from exceptions import CustomAPIException
import httpx
import asyncio
from dotenv import load_dotenv

load_dotenv()


def validate_url(url):
    domain = urlparse(url).netloc
    if domain != "www.slideshare.net":
        raise CustomAPIException(status_code=400, detail="Invalid SlideShare URL.")
    return True


def fetch_slide_images_all_resolutions(url):
    response = requests.get(url)
    if response.status_code != 200:
        raise CustomAPIException(status_code=500, detail="Failed to fetch the presentation page.")

    soup = BeautifulSoup(response.text, "html.parser")

    title_tag = soup.find("title")
    page_title = title_tag.get_text(strip=True) if title_tag else None

    image_tags = soup.find_all("img", {"data-testid": "vertical-slide-image"})

    if not image_tags:
        raise CustomAPIException(status_code=404, detail="No slide images found.")

    all_slide_images = []
    print("Beautiful soup processing")
    for tag in image_tags:
        srcset = tag.get("srcset")
        if not srcset:
            continue

        slide_resolutions = {}
        sources = [s.strip().split(" ") for s in srcset.split(",")]

        for src in sources:
            if len(src) == 2:
                url_part, res = src
                if res.endswith("w"):
                    resolution = int(res[:-1])
                    slide_resolutions[resolution] = url_part

        if slide_resolutions:
            all_slide_images.append(slide_resolutions)

    print("Beautiful soup generated images")

    return {
        "title": page_title,
        "slides": all_slide_images
    }


async def fetch_image(client: httpx.AsyncClient, url: str):
    try:
        response = await client.get(url)
        response.raise_for_status()

        img = Image.open(BytesIO(response.content)).convert("RGB")

        # Save image to temp file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        img.save(temp_file.name, format="JPEG")
        temp_file.close()  # Close so other processes can access it

        return temp_file.name  # Return the path to the temp file

    except Exception as e:
        raise CustomAPIException(
            status_code=500,
            detail=f"Failed to fetch image: {url}, Error: {str(e)}"
        )

import tempfile


semaphore = asyncio.Semaphore(10)

async def fetch_with_limit(client, url):
    async with semaphore:
        return await fetch_image(client, url)


import img2pdf


async def convert_urls_to_pdf_async(image_urls, pdf_filename):

    print('Async started')
    images = []

    # Download images asynchronously
    async with httpx.AsyncClient(timeout=20) as client:
        tasks = [fetch_with_limit(client, url) for url in image_urls]
        results = await asyncio.gather(*tasks, return_exceptions=True)

        for result in results:
            if isinstance(result, Exception):
                raise result
            images.append(result)  # Ensure all images are RGB for PDF

    if not images:
        raise CustomAPIException(status_code=500, detail="No images to convert to PDF.")

    print('Async images gathered')

    # Create temporary file
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        pdf_path = tmp_pdf.name

    try:
        # Save PDF to file
        convert_image_paths_to_pdf(images, pdf_path)

        # FTP setup
        date_str = datetime.today().strftime("%d%m%Y")
        ftp_dir = f"SS_DL/{date_str}"

        ftp_host = os.getenv("FTP_HOST")
        ftp_user = os.getenv("FTP_USER")
        ftp_pass = os.getenv("FTP_PASS")
        ftp_port = int(os.getenv("FTP_PORT", 21))

        ftp = FTP()
        ftp.set_pasv(True)
        ftp.connect(host=ftp_host, port=ftp_port)
        ftp.login(user=ftp_user, passwd=ftp_pass)
        print("Connected FTP")
        # Ensure directory exists
        for folder in ftp_dir.split('/'):
            try:
                ftp.cwd(folder)
            except:
                ftp.mkd(folder)
                ftp.cwd(folder)

        # Upload file
        with open(pdf_path, 'rb') as file_to_upload:
            ftp.storbinary(f'STOR {pdf_filename}', file_to_upload, blocksize=1048576)

        ftp.quit()
        print("File written in FTP")

        file_size = os.path.getsize(pdf_path)
        return f"{ftp_dir}/{pdf_filename}", file_size

    finally:
        os.remove(pdf_path)


def convert_image_paths_to_pdf(image_paths, pdf_path):
    try:
        with open(pdf_path, "wb") as f:
            f.write(img2pdf.convert(image_paths))
    except Exception as e:
        raise CustomAPIException(status_code=500, detail=str(e))
    finally:
        # Clean up temp image files
        for path in image_paths:
            try:
                os.remove(path)
            except Exception as cleanup_err:
                print(f"Failed to delete temp file: {path}. Error: {cleanup_err}")




from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import httpx, os, asyncio
from datetime import datetime


async def convert_urls_to_pptx_async(image_urls, pptx_filename):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    async with httpx.AsyncClient(timeout=20) as client:
        tasks = [fetch_with_limit(client, url) for url in image_urls]
        results = await asyncio.gather(*tasks, return_exceptions=True)

    for result in results:
        if isinstance(result, Exception):
            raise result

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            result, Inches(0), Inches(0),
            width=prs.slide_width, height=prs.slide_height
        )

    for result in results:
        try:
            os.remove(result)
        except Exception as e:
            print(f"Failed to delete temp file: {result}, error: {e}")

    # Save to a temporary file
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_pptx_file:
        pptx_path = tmp_pptx_file.name

    try:
        prs.save(pptx_path)

        # FTP upload
        date_str = datetime.today().strftime("%d%m%Y")
        ftp_dir = f"SS_DL/{date_str}"

        ftp_host = os.getenv("FTP_HOST")
        ftp_user = os.getenv("FTP_USER")
        ftp_pass = os.getenv("FTP_PASS")
        ftp_port = int(os.getenv("FTP_PORT", 21))

        ftp = FTP()
        ftp.set_pasv(True)
        ftp.connect(host=ftp_host, port=ftp_port)
        ftp.login(user=ftp_user, passwd=ftp_pass)

        for folder in ftp_dir.split('/'):
            try:
                ftp.cwd(folder)
            except:
                ftp.mkd(folder)
                ftp.cwd(folder)

        with open(pptx_path, 'rb') as file_to_upload:
            ftp.storbinary(f'STOR {pptx_filename}', file_to_upload,  blocksize=1048576)

        ftp.quit()

        file_size = os.path.getsize(pptx_path)
        return f"{ftp_dir}/{pptx_filename}", file_size

    finally:
        os.remove(pptx_path)


import zipfile

async def convert_urls_to_zip_async(image_urls, zip_filename):
    async with httpx.AsyncClient(timeout=20) as client:
        tasks = [fetch_with_limit(client, url) for url in image_urls]
        results = await asyncio.gather(*tasks, return_exceptions=True)

    # Create a temporary file for the zip
    with tempfile.NamedTemporaryFile(suffix=".zip", delete=False) as tmp_zip_file:
        zip_path = tmp_zip_file.name

    try:
        with zipfile.ZipFile(zip_path, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
            for idx, result in enumerate(results):
                if isinstance(result, Exception):
                    raise result
                try:
                    with open(result, "rb") as img_file:
                        zf.writestr(f"image_{idx + 1}.jpg", img_file.read())
                except Exception as e:
                    raise CustomAPIException(status_code=500, detail=str(e))
                finally:
                    # Delete the temp file after it’s written to zip
                    try:
                        os.remove(result)
                    except Exception as cleanup_err:
                        print(f"Failed to delete temp image: {result}, error: {cleanup_err}")

        # Setup FTP
        date_str = datetime.today().strftime("%d%m%Y")
        ftp_dir = f"SS_DL/{date_str}"

        ftp_host = os.getenv("FTP_HOST")
        ftp_user = os.getenv("FTP_USER")
        ftp_pass = os.getenv("FTP_PASS")
        ftp_port = int(os.getenv("FTP_PORT", 21))

        ftp = FTP()
        ftp.set_pasv(True)
        ftp.connect(host=ftp_host, port=ftp_port)
        ftp.login(user=ftp_user, passwd=ftp_pass)

        for folder in ftp_dir.split('/'):
            try:
                ftp.cwd(folder)
            except:
                ftp.mkd(folder)
                ftp.cwd(folder)

        # Upload from temp file
        with open(zip_path, 'rb') as file_to_upload:
            ftp.storbinary(f'STOR {zip_filename}', file_to_upload,  blocksize=1048576)

        ftp.quit()

        # Get final file size before deleting
        file_size = os.path.getsize(zip_path)

        return f"{ftp_dir}/{zip_filename}", file_size

    finally:
        # Always delete temp file
        os.remove(zip_path)


from utils import SlidesConversionType, QualityType
from ftplib import FTP


async def get_slides_download_link(url: str, conversion_type: SlidesConversionType,
                                   quality_type: QualityType = QualityType.hd):
    validate_url(url)

    path_parts = urlparse(url).path.strip("/").split("/")
    if len(path_parts) >= 2:
        doc_short = path_parts[-2]
    else:
        raise CustomAPIException(status_code=400, detail="Invalid SlideShare URL format.")

    fetch_slide_images_title = fetch_slide_images_all_resolutions(url)

    slide_images = fetch_slide_images_title['slides']
    title = fetch_slide_images_title['title']



    quality = 2048 if quality_type == QualityType.hd else 638

    high_res_images = [slide[quality] for slide in slide_images if quality in slide]

    if not high_res_images:
        raise CustomAPIException(status_code=404, detail=f"No {quality}px resolution slides found.")

    thumbnail = high_res_images[0]


    if conversion_type == SlidesConversionType.pdf:
        path, total_size = await convert_urls_to_pdf_async(high_res_images, f"{doc_short}.pdf")
        message = "PDF generated successfully."
    elif conversion_type == SlidesConversionType.pptx:
        path, total_size = await convert_urls_to_pptx_async(high_res_images, f"{doc_short}.pptx")
        message = "PPTX generated successfully."
    elif conversion_type == SlidesConversionType.images_zip:
        path, total_size = await convert_urls_to_zip_async(high_res_images, f"{doc_short}.zip")
        message = "IMAGES ZIP generated successfully."
    else:
        raise CustomAPIException(status_code=400, detail="Unsupported conversion type.")

    file_name = os.path.basename(path)

    base_url = os.getenv("BASE_URL", 21)

    return {
        "success": True,
        "message": message,
        "data": json.dumps({
            "thumbnail": thumbnail,
            "quality": quality_type.value,
            "conversion_type": conversion_type.value,
            "slides_download_link": f'{base_url}/{path}',
            "file_name": file_name,
            "size": total_size,
            "title": title
        })

    }
